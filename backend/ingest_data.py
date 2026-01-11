import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from sqlmodel import Session, create_engine, select, SQLModel
from app.models import KnowledgeItem 

# إعداد الاتصال
sqlite_url = "sqlite:///database.db" 
engine = create_engine(sqlite_url)

def setup_database():
    print("🛠️ جاري تهيئة قاعدة البيانات...")
    SQLModel.metadata.create_all(engine)

def extract_text_from_pdf(path):
    text = ""
    try:
        with fitz.open(path) as doc:
            for page in doc: text += page.get_text()
    except Exception as e: print(f"   ⚠️ خطأ PDF: {e}")
    return text

def extract_text_from_word(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs: text += para.text + "\n"
    except Exception as e: print(f"   ⚠️ خطأ Word: {e}")
    return text

def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e: print(f"   ⚠️ خطأ PPTX: {e}")
    return text

def start_ingestion():
    folders = [
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\iq",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L2 Grade 10",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L3 Grade 11",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L3 Grade 12"
    ]

    setup_database()

    with Session(engine) as session:
        print("\n--- 🚀 سحب البيانات الشامل (PDF, Word, PPTX) ---")
        
        for folder in folders:
            if not os.path.exists(folder): continue
            folder_name = os.path.basename(folder)
            
            for filename in os.listdir(folder):
                path = os.path.join(folder, filename)
                ext = filename.lower()
                
                # تخطي إذا كان مضافاً مسبقاً
                existing = session.exec(select(KnowledgeItem).where(KnowledgeItem.source_file == filename)).first()
                if existing: continue

                content = ""
                if ext.endswith(".pdf"):
                    print(f"   📄 معالجة PDF: {filename}...")
                    content = extract_text_from_pdf(path)
                elif ext.endswith(".docx"):
                    print(f"   📝 معالجة Word: {filename}...")
                    content = extract_text_from_word(path)
                elif ext.endswith(".pptx"):
                    print(f"   📊 معالجة PowerPoint: {filename}...")
                    content = extract_text_from_pptx(path)

                if content.strip():
                    item = KnowledgeItem(content=content, source_file=filename, category=folder_name)
                    session.add(item)
        
        session.commit()
        print("\n✨ مبروك! المساعد الآن خبير في الكتب والأبحاث والعروض التقديمية.")

if __name__ == "__main__":
    start_ingestion()