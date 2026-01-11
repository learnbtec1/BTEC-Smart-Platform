import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from sqlmodel import Session, create_engine, select, SQLModel
from app.models import KnowledgeItem 

# إعداد الاتصال
sqlite_url = "sqlite:///database.db" 
engine = create_engine(sqlite_url)

def setup_database():
    print("🛠️ جاري تهيئة قاعدة البيانات وتحديث الهيكل...")
    SQLModel.metadata.create_all(engine)

def extract_text_from_pdf_with_pages(path):
    """استخراج النص مع حفظ رقم الصفحة لكل جزء"""
    pages_content = []
    try:
        with fitz.open(path) as doc:
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                if text.strip():
                    pages_content.append({"page": page_num, "text": text})
    except Exception as e:
        print(f"   ⚠️ خطأ PDF في {os.path.basename(path)}: {e}")
    return pages_content

def extract_text_from_word(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs: text += para.text + "\n"
    except Exception as e: print(f"   ⚠️ خطأ Word: {e}")
    return text

def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text += shape.text + "\n"
    except Exception as e: print(f"   ⚠️ خطأ PPTX: {e}")
    return text

def start_ingestion():
    folders = [
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\iq",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L2 Grade 10",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L3 Grade 11",
        r"D:\BTEC-backend\backend\app\knowledge_base\Business\L3 Grade 12"
    ]

    setup_database()

    with Session(engine) as session:
        print("\n--- 🚀 بدء سحب البيانات مع دعم أرقام الصفحات ---")
        
        for folder in folders:
            if not os.path.exists(folder): continue
            folder_name = os.path.basename(folder)
            
            for filename in os.listdir(folder):
                path = os.path.join(folder, filename)
                ext = filename.lower()
                
                # تخطي الملف إذا تم سحبه مسبقاً (اختياري: يمكنك مسح قاعدة البيانات للبدء من جديد مع أرقام الصفحات)
                existing = session.exec(select(KnowledgeItem).where(KnowledgeItem.source_file == filename)).first()
                if existing: continue

                if ext.endswith(".pdf"):
                    print(f"   📄 معالجة PDF مع الصفحات: {filename}...")
                    pages = extract_text_from_pdf_with_pages(path)
                    for p in pages:
                        item = KnowledgeItem(
                            content=p["text"],
                            source_file=filename,
                            category=f"{folder_name} (Page {p['page']})" # تخزين رقم الصفحة في التصنيف أو النص
                        )
                        session.add(item)
                
                elif ext.endswith((".docx", ".pptx")):
                    print(f"   📝 معالجة ملف مكتب: {filename}...")
                    content = extract_text_from_word(path) if ext.endswith(".docx") else extract_text_from_pptx(path)
                    if content.strip():
                        item = KnowledgeItem(content=content, source_file=filename, category=folder_name)
                        session.add(item)
        
        session.commit()
        print("\n✨ تم سحب البيانات بنجاح مع أرقام الصفحات!")

if __name__ == "__main__":
    start_ingestion()